"""Fills templates/guide_request_template.pptx (the real ガイド依頼書 master)
with a booking's data.

Unlike the xlsx template, python-pptx round-trips this file cleanly (text
runs are preserved), so we can use its normal API instead of raw XML
surgery. See docgen/xlsx_lowlevel.py for why the xlsx template needs a
different approach.

Each source paragraph is a scatter of many same-content runs (an artifact
of how Google Slides saves edits), with no run boundary lining up with the
label/value split we need. So instead of tweaking individual runs, we
collapse each paragraph down to a single run and set its whole text --
this keeps paragraph-level formatting (alignment, indent) while letting us
control font.bold explicitly where the business rule requires it (the
dietary/allergy warning must be bold Japanese).
"""

from __future__ import annotations

import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from ..models import BookingRequest, ItineraryStop
from ..rules import format_dietary_for_guide_request, rewrite_payment_label_for_pax, tba, tbd

TEMPLATE_PATH = (
    Path(__file__).resolve().parent.parent.parent.parent / "templates" / "guide_request_template.pptx"
)

ITINERARY_HEADER_ROWS = 1
ITINERARY_MAX_ROWS = 17  # template table has 18 rows total, minus the header
DEFAULT_TABLE_FONT_SIZE = Pt(9)  # matches the original template's itinerary rows


def _set_paragraph_text(
    paragraph, text: str, *, bold: bool | None = None, size=None, font_name: str | None = None
) -> None:
    runs = paragraph.runs
    if not runs:
        run = paragraph.add_run()
    else:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    run.text = text
    if bold is not None:
        run.font.bold = bold
    if size is not None:
        run.font.size = size
    if font_name is not None:
        run.font.name = font_name


def _participants_text(booking: BookingRequest) -> str:
    return ", ".join(
        f"{p.name} ({p.age})" if p.age is not None else p.name for p in booking.participants
    )


def _display_date_jp(iso_date: str) -> str:
    try:
        parsed = date.fromisoformat(iso_date)
    except ValueError:
        return iso_date
    weekday = "月火水木金土日"[parsed.weekday()]
    return f"{parsed.year}年{parsed.month}月{parsed.day}日（{weekday}）"


def _fee_text(booking: BookingRequest) -> str:
    """Total guide fee, plus a compact breakdown when there's a non-zero
    adjustment (e.g. extra hotel-pickup time) so the guide/office can see
    what the total is actually made of, not just the final figure."""
    if booking.guide_fee is None:
        return tbd(None)
    text = f"{booking.guide_fee:,}円"
    if booking.guide_fee_adjustment:
        base = booking.guide_fee - booking.guide_fee_adjustment
        text += f"（基本{base:,}／調整{booking.guide_fee_adjustment:+,}）"
    return text


def _find_shape(slide, contains: str):
    for shape in slide.shapes:
        if shape.has_text_frame and contains in shape.text_frame.text:
            return shape
    raise ValueError(f"no shape containing {contains!r} found")


def _fill_header(slide, booking: BookingRequest) -> None:
    guide_shape = _find_shape(slide, "様")
    _set_paragraph_text(guide_shape.text_frame.paragraphs[0], f" {booking.guide_name}　様　　")

    issue_shape = _find_shape(slide, "発行日")
    today = date.today().isoformat()
    _set_paragraph_text(
        issue_shape.text_frame.paragraphs[0], f"発行日：　          {today}"
    )

    company_shape = _find_shape(slide, "担当者")
    _set_paragraph_text(
        company_shape.text_frame.paragraphs[4], f"担当者 {booking.assignee_1st or tba(None)}"
    )

    details_shape = _find_shape(slide, "ツアー名")
    paragraphs = details_shape.text_frame.paragraphs
    _set_paragraph_text(paragraphs[0], f"日程　　　　{_display_date_jp(booking.tour_date)}")
    _set_paragraph_text(paragraphs[1], f"ツアー名 \t{booking.tour_name}")
    _set_paragraph_text(paragraphs[2], f"集合場所\t{booking.meeting_point_jp or tba(None)}")
    _set_paragraph_text(paragraphs[3], f"ゲスト名 \t{_participants_text(booking) or tba(None)}")
    _set_paragraph_text(paragraphs[4], f"時間 \t{booking.start_time} - {booking.end_time}")
    _set_paragraph_text(paragraphs[5], f"人数 \t{booking.pax}名")
    fee_text = _fee_text(booking)
    _set_paragraph_text(paragraphs[6], f"ガイド謝金  \t{fee_text}  　　　　　　　【行程】")


def _existing_font(text_frame):
    """Grab the first run's font size/name before we clear the cell.

    `TextFrame.clear()` drops all runs, and a freshly added run falls back
    to PowerPoint's default (much larger) size -- without this, a fully
    populated itinerary overflows the slide even though the same text fit
    fine in the original template.
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return run.font.size, run.font.name
    return DEFAULT_TABLE_FONT_SIZE, "Arial"


def _fill_itinerary_table(slide, booking: BookingRequest) -> None:
    table_shape = next(shape for shape in slide.shapes if shape.has_table)
    table = table_shape.table
    stops: list[ItineraryStop] = booking.itinerary[:ITINERARY_MAX_ROWS]
    for row_offset in range(ITINERARY_MAX_ROWS):
        row = table.rows[ITINERARY_HEADER_ROWS + row_offset]
        stop = stops[row_offset] if row_offset < len(stops) else None
        values = (
            (
                stop.time_label,
                stop.stopover_name,
                rewrite_payment_label_for_pax(stop.payment_label, booking.pax),
                stop.payment_method,
                stop.stopover_info,
            )
            if stop
            else ("", "", "", "", "")
        )
        for cell, value in zip(row.cells, values):
            text_frame = cell.text_frame
            size, font_name = _existing_font(text_frame)
            text_frame.clear()  # drop any extra leftover paragraphs from the template row
            _set_paragraph_text(text_frame.paragraphs[0], value, size=size, font_name=font_name)


def _fill_notes(slide, booking: BookingRequest) -> None:
    notes_shape = _find_shape(slide, "注意事項")
    paragraphs = notes_shape.text_frame.paragraphs

    dietary_jp = format_dietary_for_guide_request(booking.dietary)
    _set_paragraph_text(paragraphs[0], f"【注意事項】{dietary_jp}", bold=True)

    _set_paragraph_text(
        paragraphs[1], f"【緊急連絡先】{booking.emergency_contact or tbd(None)}"
    )


def generate_guide_request(booking: BookingRequest, output_path: Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE_PATH, output_path)

    prs = Presentation(output_path)
    slide = prs.slides[0]

    _fill_header(slide, booking)
    _fill_itinerary_table(slide, booking)
    _fill_notes(slide, booking)

    prs.save(output_path)
    return output_path
